import collections 
import collections.abc
from pptx import Presentation
import pandas as pd
import re

def clean_text(text):
    if not text:
        return ""
    return text.replace('\xa0', ' ').strip()

def extract_score(text):
    match = re.search(r"(\d+(\.\d+)?)", str(text))
    if match:
        return float(match.group(1))
    return 0.0

def extract_percent(text):
    match = re.search(r"(\d+(\.\d+)?)", str(text))
    if match:
        return float(match.group(1))
    return 0.0

def main():
    prs = Presentation("Assessment Analyze.pptx")
    data = []

    print("Extracting data from slides...")
    
    active_record = None
    current_team = "Unassigned"
    
    for i, slide in enumerate(prs.slides):
        # 1. Detect Team Section Header
        slide_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text += shape.text + " "
        
        slide_text_lower = slide_text.lower().strip()
        
        if "qa/qc" in slide_text_lower and len(slide_text) < 100:
            current_team = "QA/QC"
        elif "r&d" in slide_text_lower and len(slide_text) < 100:
            current_team = "R&D"
        elif ("batt dev" in slide_text_lower or "battery development" in slide_text_lower) and len(slide_text) < 100:
            current_team = "Battery Development"
        elif "production" in slide_text_lower and len(slide_text) < 50:
            current_team = "Production"

        # 2. Extract Table Data
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            
            table = shape.table
            if len(table.rows) < 2: continue
            
            header_cells = [c.text_frame.text.lower().strip() for c in table.rows[0].cells]
            
            # Helper to find headers
            score_col_idx = -1
            notes_col_idx = -1
            
            for idx, h in enumerate(header_cells):
                if any(k in h for k in ["performance", "quantification", "rating"]):
                     score_col_idx = idx
                     break
            
            if score_col_idx == -1:
                for idx, h in enumerate(header_cells):
                    if "score" in h and "reason" not in h:
                         score_col_idx = idx
                    if "weight" in h:
                         score_col_idx = idx

            for idx, h in enumerate(header_cells):
                if "comments" in h or "notes" in h or "reason" in h:
                    notes_col_idx = idx

            # Check if this table has a Name Row
            has_name = False
            for row in table.rows:
                if len(row.cells) > 0 and "employee name" in row.cells[0].text_frame.text.lower():
                    has_name = True
                    break
            
            # State Management
            if has_name:
                # Flush previous
                if active_record and 'Employee Name' in active_record:
                    if active_record['Employee Name'].strip() and active_record['Employee Name'] != "Unassigned":
                        data.append(active_record)
                # Start new
                active_record = {'Team': current_team}
            elif active_record is None:
                # Orphan table with no preceding name
                continue
                
            # Parse Rows into active_record
            for row in table.rows:
                cells = [clean_text(cell.text_frame.text) for cell in row.cells]
                if not cells: continue
                
                header = cells[0].lower()
                
                # --- Helpers ---
                def get_score(cells, idx):
                    if idx != -1 and idx < len(cells):
                         return extract_score(cells[idx])
                    if len(cells) == 2: return extract_score(cells[1])
                    if len(cells) > 2 and "%" in cells[1]: return extract_score(cells[2])
                    return 0.0

                def get_notes(cells, idx):
                    if idx != -1 and idx < len(cells):
                        return cells[idx]
                    if len(cells) > 4: return cells[4]
                    if len(cells) > 2 and idx == -1: return cells[-1]
                    return ""
                
                # --- Extraction ---
                if "employee name" in header:
                    c0 = cells[0].lower()
                    if "(employee name)" in c0:
                        val = re.sub(r"\(employee name\)", "", cells[0], flags=re.IGNORECASE).strip()
                        active_record['Employee Name'] = val
                    elif c0 == "employee name" and len(cells) > 1:
                        active_record['Employee Name'] = cells[1].strip()
                    else:
                        active_record['Employee Name'] = cells[0].strip()
                
                elif "role" in header:
                    c0 = cells[0].lower()
                    if "(role)" in c0:
                        val = re.sub(r"\(role\)", "", cells[0], flags=re.IGNORECASE).strip()
                        active_record['Role'] = val
                    elif c0 == "role" and len(cells) > 1:
                        active_record['Role'] = cells[1].strip()
                    else:
                        active_record['Role'] = cells[0].strip()
                
                elif "quality" in header:
                    active_record['Quality Score'] = get_score(cells, score_col_idx)
                    active_record['Quality Notes'] = get_notes(cells, notes_col_idx)
                elif "productivity" in header:
                    active_record['Productivity Score'] = get_score(cells, score_col_idx)
                    active_record['Productivity Notes'] = get_notes(cells, notes_col_idx)
                elif "attendance" in header:
                    active_record['Attendance Score'] = get_score(cells, score_col_idx)
                    active_record['Attendance Notes'] = get_notes(cells, notes_col_idx)
                elif "skill" in header:
                    active_record['Skill Score'] = get_score(cells, score_col_idx)
                    active_record['Skill Notes'] = get_notes(cells, notes_col_idx)
                elif "teamwork" in header:
                    active_record['Teamwork Score'] = get_score(cells, score_col_idx)
                    active_record['Teamwork Notes'] = get_notes(cells, notes_col_idx)

                elif "weighted final score" in header:
                    active_record['Weighted Score'] = get_score(cells, score_col_idx)

                elif "recommended salary increase" in header:
                     val = get_score(cells, score_col_idx)
                     if val == 0.0 and len(cells) > 1:
                         val = extract_percent(cells[1])
                     active_record['Salary Increase'] = val

                elif "manager comments" in header or "notes" in header:
                     content = " ".join(cells[1:])
                     active_record['Manager Notes'] = content.strip()

    # Flush last record
    if active_record and 'Employee Name' in active_record:
        if active_record['Employee Name'].strip() and active_record['Employee Name'] != "Unassigned":
            data.append(active_record)

    df = pd.DataFrame(data)
    # Deduplicate
    df = df.drop_duplicates(subset=['Employee Name'])
    
    print(f"Extracted {len(df)} records.")
    
    # Save
    df.to_csv("employee_data.csv", index=False)
    print(df[['Employee Name', 'Team', 'Weighted Score', 'Salary Increase']])

if __name__ == "__main__":
    main()
